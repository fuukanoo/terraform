import logging
from pptx import Presentation
from io import BytesIO

def extract_text_from_ppt(blob_data):
    """PowerPointファイルのバイナリデータからテキストを抽出"""
    try:
        # BytesIOを使用してメモリ上でPowerPointファイルを処理
        prs = Presentation(BytesIO(blob_data))
        
        # すべてのスライドのテキストを抽出
        extracted_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
        
        return extracted_text

    except Exception as e:
        logging.error(f"Error processing the PowerPoint document: {e}")
        return f"Error processing the PowerPoint document: {e}"