"""buckup for function_app.py"""

import azure.functions as func
import os
import openai
from azure.storage.blob import BlobServiceClient
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
import requests
import logging

import mimetypes  # get file type
from docx import Document
import pandas as pd
from pptx import Presentation
from io import BytesIO

# get API_KEY and ENDPOINT
vision_key = os.getenv("COMPUTER_VISION_API_KEY")
vision_endpoint = os.getenv("COMPUTER_VISION_ENDPOINT")
intelligence_key = os.getenv("DOCUMENT_INTELLIGENCE_API_KEY")
intelligence_endpoint = os.getenv("DOCUMENT_INTELLIGENCE_ENDPOINT")
blob_connection_string = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
blob_container = os.getenv("BLOB_CONTAINER")
openai.api_key = os.getenv("AZURE_OPENAI_API_KEY")
openai.azure_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")  
openai.api_type = "azure"
openai.api_version = "2023-05-15"  # version

# route setting
app = func.FunctionApp(http_auth_level=func.AuthLevel.ANONYMOUS)

@app.route(route="function_rag")
def main(req: func.HttpRequest) -> func.HttpResponse:
    logging.info('Processing a request to extract and improve text from an image or document.')

    # リクエストボディからJSONデータを取得
    try:
        req_body = req.get_json()
    except ValueError:
        return func.HttpResponse("Invalid JSON format", status_code=400)

    file_name = req_body.get("file_name")

    if not file_name:
        return func.HttpResponse("Please provide 'file_name' in the request body.", status_code=400)

    # Blobストレージからファイルを取得
    extracted_text = extract_text_from_blob(blob_container, file_name)

    if "Error" in extracted_text:
        return func.HttpResponse(extracted_text, status_code=500)

    improved_text = improve_text_with_openai(extracted_text)
    logging.info(f"Improved text:\n {improved_text}")

    return func.HttpResponse(improved_text, status_code=200)


def extract_text_from_blob(blob_container, blob_name):
    """ファイルの種類に応じてBlobストレージから画像、ドキュメント、Word、Excel、PowerPointを取得し、適切な方法でテキストを抽出"""
    blob_service_client = BlobServiceClient.from_connection_string(blob_connection_string)
    blob_client = blob_service_client.get_blob_client(container=blob_container, blob=blob_name)

    try:
        # Blobからファイルをダウンロード
        blob_data = blob_client.download_blob().readall()

        # MIMEタイプを取得して処理を分岐
        mime_type, _ = mimetypes.guess_type(blob_name)

        if mime_type and mime_type.startswith("image"):
            return extract_text_from_image(blob_data)
        elif mime_type == "application/pdf":
            return extract_text_from_document(blob_data)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            #ipdb.set_trace()
            return extract_text_from_word(blob_data)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return extract_text_from_excel(blob_data)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_text_from_ppt(blob_data)
        else:
            return "Error: Unsupported file type."

    except Exception as e:
        logging.error(f"Error processing the blob: {e}")
        return f"Error processing the blob: {e}"


def extract_text_from_image(blob_data):
    """Computer Vision APIを使って画像データから文字を抽出する"""
    ocr_url = vision_endpoint + "/vision/v3.1/ocr"
    headers = {
        'Ocp-Apim-Subscription-Key': vision_key,
        'Content-Type': 'application/octet-stream'
    }

    try:
        response = requests.post(ocr_url, headers=headers, data=blob_data)
        response.raise_for_status()
        analysis = response.json()

        # OCR結果からテキストを抽出
        extracted_text = ""
        for region in analysis.get("regions", []):
            for line in region.get("lines", []):
                line_text = " ".join([word["text"] for word in line.get("words", [])])
                extracted_text += line_text + "\n"

        return extracted_text

    except requests.exceptions.HTTPError as e:
        logging.error(f"Error processing the image: {e}")
        return f"Error processing the image: {e}"
    except Exception as e:
        logging.error(f"An error occurred while extracting text from the image: {e}")
        return f"An error occurred while extracting text from the image: {e}"


def extract_text_from_document(blob_data):
    """Document Intelligence APIを使ってドキュメントから文字を抽出する"""
    document_analysis_client = DocumentAnalysisClient(
        endpoint=intelligence_endpoint, 
        credential=AzureKeyCredential(intelligence_key)
    )

    try:
        # Document Intelligence API に送信
        poller = document_analysis_client.begin_analyze_document("prebuilt-document", blob_data)
        result = poller.result()

        # 結果からテキストを抽出
        extracted_text = ""
        for page in result.pages:
            for line in page.lines:
                extracted_text += line.content + "\n"
        return extracted_text

    except Exception as e:
        logging.error(f"Error processing the document: {e}")
        return f"Error processing the document: {e}"
    

def extract_text_from_word(blob_data):
    """Wordファイルのバイナリデータからテキストを抽出"""
    try:
        # BytesIOを使ってバイナリデータをメモリ上で処理
        doc = Document(BytesIO(blob_data))
        
        # 文書内のすべての段落を抽出し、改行で結合
        extracted_text = "\n".join([para.text for para in doc.paragraphs])
        
        return extracted_text

    except Exception as e:
        logging.error(f"Error processing the Word document: {e}")
        return f"Error processing the Word document: {e}"


def extract_text_from_excel(blob_data):
    """Excelファイルのバイナリデータからテキストを抽出"""
    try:
        # BytesIOを使用してメモリ上でExcelファイルを処理
        df = pd.read_excel(BytesIO(blob_data))
        
        # DataFrameを文字列形式に変換
        extracted_text = df.to_string(index=False)
        
        return extracted_text

    except Exception as e:
        logging.error(f"Error processing the Excel document: {e}")
        return f"Error processing the Excel document: {e}"


def extract_text_from_ppt(blob_data):
    """PowerPointファイルのバイナリデータからテキストを抽出"""
    try:
        # BytesIOを使用してメモリ上でPowerPointファイルを処理
        prs = Presentation(BytesIO(blob_data))
        
        # すべてのスライドのテキストを抽出
        extracted_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
        
        return extracted_text

    except Exception as e:
        logging.error(f"Error processing the PowerPoint document: {e}")
        return f"Error processing the PowerPoint document: {e}"


def improve_text_with_openai(text):
    """Azure OpenAIを使って抽出したテキストを修正する"""
    prompt = (f"次のテキストはOCRを使用して画像から,もしくはドキュメントから抽出されました. "
              f"誤読された部分や間違っている部分があれば，訂正して正しいと思われる文章にして出力して. "
              f"ただし,正しい場合は特に変更を加えずにそのまま出力してください.:\n\n{text}")

    try:
        # Azure OpenAI API を使ってテキスト補正
        response = openai.chat.completions.create(
            model="gpt-35-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant "},
                {"role": "user", "content": prompt},
            ],
            max_tokens=1000,
            temperature=0.3,
            top_p=0.9
        )

        message_content = response.choices[0].message.content.strip()
        return message_content

    except openai.OpenAIError as e:
        logging.error(f"Error during OpenAI request: {e}")
        return f"Error during OpenAI request: {e}"